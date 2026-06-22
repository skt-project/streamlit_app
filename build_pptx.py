"""Generates User_Training.pptx from the screenshots captured by
docs_crawler.py, structured into the 9 required sections.

Usage: python build_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent
SHOTS = ROOT / "screenshots"

NAVY = RGBColor(0x14, 0x50, 0x6A)
TEAL = RGBColor(0x1E, 0x6B, 0x8A)
LIGHT = RGBColor(0xEA, 0xF6, 0xFB)
DARK_TEXT = RGBColor(0x1C, 0x1C, 0x1C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color=RGBColor(0xFF, 0xFF, 0xFF)):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.color.rgb = RGBColor(0xCC, 0xE8, 0xF0)


def add_bullets(slide, bullets, left=0.6, top=1.4, width=5.6, height=5.6, font_size=16):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = ("• " if level == 0 else "‣ ") + text
        run.font.size = Pt(font_size if level == 0 else font_size - 2)
        run.font.color.rgb = DARK_TEXT
        run.font.bold = level == 0 and text.endswith(":")
        p.space_after = Pt(6)
    return tb


def add_image_right(slide, filename, left=6.4, top=1.4, width=6.3, max_height=5.8):
    path = SHOTS / filename
    if path.exists():
        from PIL import Image
        with Image.open(path) as im:
            w, h = im.size
        if (width * (h / w)) > max_height:
            height = max_height
            width = height * (w / h)
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1))
        tb.text_frame.text = f"[missing screenshot: {filename}]"


def add_image_full(slide, filename, top=1.3, max_width=12.1, max_height=5.9):
    path = SHOTS / filename
    if not path.exists():
        return
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    ratio = min(max_width / (w / 96), max_height / (h / 96))
    disp_w = Inches((w / 96) * ratio)
    disp_h = Inches((h / 96) * ratio)
    left = Inches((13.333 - disp_w.inches) / 2)
    slide.shapes.add_picture(str(path), left, Inches(top), width=disp_w, height=disp_h)


def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(12), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x90, 0x90, 0x90)


# ============================================================
# SECTION 1 — INTRODUCTION
# ============================================================
slide = add_slide()
add_bg(slide, NAVY)
tb = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.3), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Distributor Operational Assessment"
run.font.size = Pt(44)
run.font.bold = True
run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "User Training — Area Sales Supervisor Role"
run2.font.size = Pt(22)
run2.font.color.rgb = RGBColor(0xBE, 0xE3, 0xEE)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "SKINTIFIC"
run3.font.size = Pt(16)
run3.font.bold = True
run3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

slide = add_slide()
add_title_bar(slide, "Introduction", "What this app does and who it's for")
add_bullets(slide, [
    ("What it is:", 0),
    ("A multi-role scoring tool that assesses each distributor across 10 operational metrics every month.", 1),
    ("Who's involved:", 0),
    ("Area Sales Supervisor — people, infrastructure, delivery & inventory (7 metrics, 34 pts)", 1),
    ("Distributor Manager — bulk allocation uploads (NPD, SKU Focus) + user admin", 1),
    ("Admin RSA — operational scoring inputs", 1),
    ("Account Receivable — AR performance + data reporting compliance", 1),
    ("This guide covers:", 0),
    ("The Area Sales Supervisor workflow end-to-end, from login to submission.", 1),
], width=12.2)
add_footer(slide, "Distributor Operational Assessment — User Training")

# ============================================================
# SECTION 2 — LOGIN
# ============================================================
slide = add_slide()
add_title_bar(slide, "2. Login", "Authenticate and reach your role's dashboard")
add_bullets(slide, [
    ("Purpose:", 0),
    ("Your account determines what you see — there is no manual role picker.", 1),
    ("Fields:", 0),
    ("Username (not case-sensitive)", 1),
    ("Password", 1),
    ("Validation:", 0),
    ("Wrong credentials show one generic error — doesn't reveal which field was wrong.", 1),
    ("No retry/lockout limit is enforced.", 1),
], width=6.0)
add_image_right(slide, "01_login_screen.png")
add_footer(slide, "Distributor Operational Assessment — User Training")

# ============================================================
# SECTION 3 — DASHBOARD OVERVIEW
# ============================================================
slide = add_slide()
add_title_bar(slide, "3. Dashboard Overview", "What you see right after logging in")
add_bullets(slide, [
    ("Sidebar:", 0),
    ("Your name, role, region, Logout, Change Password", 1),
    ("Live Score and per-metric breakdown (updates as you fill the form)", 1),
    ("Assessment Progress — all 4 roles' status for this distributor+period", 1),
    ("Main area:", 0),
    ("Logged In As card — confirms identity & scope", 1),
    ("Assessment Period selector", 1),
    ("Location: Region (locked) + Distributor (required to unlock the form)", 1),
], width=6.0)
add_image_right(slide, "03_post_login_locked.png")
add_footer(slide, "Distributor Operational Assessment — User Training")

# ============================================================
# SECTION 4 — ASSESSMENT WORKFLOW
# ============================================================
slide = add_slide()
add_title_bar(slide, "4. Assessment Workflow", "Step 1 — choose period and distributor")
add_bullets(slide, [
    ("Assessment Period:", 0),
    ("Area Sales Supervisor sees only the current month + 2 prior months.", 1),
    ("Distributor:", 0),
    ("List is filtered to distributors mapped to YOU specifically, not the whole region.", 1),
    ("Empty list = nothing mapped to your account yet — contact an Admin.", 1),
], width=6.0)
add_image_right(slide, "05_distributor_dropdown.png")
add_footer(slide, "Distributor Operational Assessment — User Training")

slide = add_slide()
add_title_bar(slide, "4. Assessment Workflow", "Step 2 — People & Roles (metrics 1–3)")
add_bullets(slide, [
    ("Operational Leader — max 5 pts", 0),
    ("Salesman — max 7 pts, requires exactly 5 names", 0),
    ("Administrative & AR Support — max 3 pts", 0),
    ("Validation:", 0),
    ("Grade ≠ 'Do not exist' → a name is required.", 1),
    ("Grade = 'Do not exist' → no name should be entered.", 1),
], width=6.0)
add_image_right(slide, "07_people_roles_filled.png")
add_footer(slide, "Distributor Operational Assessment — User Training")

slide = add_slide()
add_title_bar(slide, "4. Assessment Workflow", "Step 3 — Infrastructure, Delivery & Operations (metrics 4–6, 9)")
add_bullets(slide, [
    ("Warehouse Facility Standard — max 3 pts (direct grade)", 0),
    ("Delivery SLA Compliance — max 8 pts (calculated)", 0),
    ("Inner/Outer City bands; both at 100% = 8 pts, any <80% = 0 pts", 1),
    ("Inventory Control & Stock Opname — max 6 pts (direct grade)", 0),
    ("Bad Stock Handling Performance — max 2 pts (calculated)", 0),
    ("Compliance % = Utilization ÷ (0.5% of YTD Sell Through) × 100", 1),
    ("Card hides + auto-scores max if YTD Sell Through is zero/unavailable", 1),
], width=6.0, font_size=15)
add_image_right(slide, "09_operations_compliance.png")
add_footer(slide, "Distributor Operational Assessment — User Training")

# ============================================================
# SECTION 5 — SCORE SUBMISSION
# ============================================================
slide = add_slide()
add_title_bar(slide, "5. Score Submission", "Live tracking, review, and final confirmation")
add_bullets(slide, [
    ("Live Score Tracker (sidebar):", 0),
    ("Updates instantly as you answer — running total out of 34.", 1),
    ("Shows other roles' progress on the same distributor+period.", 1),
    ("Review & Submit:", 0),
    ("Validates the form, then opens a confirmation popup with the full breakdown.", 1),
    ("One submission per role, per distributor, per assessment period — final once confirmed.", 1),
], width=6.0)
add_image_right(slide, "12_review_submit_popup.png", width=6.5)
add_footer(slide, "Distributor Operational Assessment — User Training")

# ============================================================
# SECTION 6 — DISTRIBUTOR MAPPING
# ============================================================
slide = add_slide()
add_title_bar(slide, "6. Distributor Mapping", "Why you only see certain distributors")
add_bullets(slide, [
    ("How the list is built:", 0),
    ("Matched against the master distributor list using YOUR region and full name", 1),
    ("Specifically: your name must match the assigned supervisor field for SKINTIFIC or Timephoria on that distributor", 1),
    ("Why this matters:", 0),
    ("Prevents assessing distributors that aren't actually yours — keeps accountability clean.", 1),
    ("If something's missing:", 0),
    ("Contact an Admin to check/update the master distributor mapping — this app doesn't let you self-add distributors.", 1),
], width=12.2, font_size=16)
add_footer(slide, "Distributor Operational Assessment — User Training")

# ============================================================
# SECTION 7 — REPORTS
# ============================================================
slide = add_slide()
add_title_bar(slide, "7. Reports", "Where to see scores today (and what's not built yet)")
add_bullets(slide, [
    ("Available now:", 0),
    ("Sidebar 'Assessment Progress' — live status of all 10 metrics across all 4 roles for the distributor+period you're viewing.", 1),
    ("'Combined Score So Far' — running total out of 50 as roles submit.", 1),
    ("Confirmation popup — your own role's full point breakdown at submission time.", 1),
    ("Not available yet:", 0),
    ("There is no dedicated multi-distributor or historical reporting screen inside this app — that would be a future enhancement, not a current feature.", 1),
], width=7.6, font_size=16)
add_image_right(slide, "11_sidebar_score_tracker.png", left=8.4, top=1.3, width=4.3)
add_footer(slide, "Distributor Operational Assessment — User Training")

# ============================================================
# SECTION 8 — FAQ
# ============================================================
slide = add_slide()
add_title_bar(slide, "8. FAQ", "")
add_bullets(slide, [
    ("Q: I can't find a distributor I expect to see.", 0),
    ("A: Your list is filtered to distributors mapped to you. Contact an Admin if one's missing.", 1),
    ("Q: The Bad Stock card disappeared — bug?", 0),
    ("A: No — happens when YTD sell-through is zero. Full marks are auto-awarded.", 1),
    ("Q: I made a mistake after confirming — can I fix it?", 0),
    ("A: No. Submissions are final per role/distributor/period. Contact an Admin.", 1),
    ("Q: Why does the sidebar show other roles' progress?", 0),
    ("A: All 4 roles roll up into one combined assessment per distributor+period.", 1),
], width=12.2, font_size=16)
add_footer(slide, "Distributor Operational Assessment — User Training")

# ============================================================
# SECTION 9 — TROUBLESHOOTING
# ============================================================
slide = add_slide()
add_title_bar(slide, "9. Troubleshooting", "")
rows = [
    ("Symptom", "Likely Cause", "Fix"),
    ("\"Invalid username or password\"", "Typo or wrong account", "Re-check credentials (password is case-sensitive)"),
    ("Distributor dropdown is empty", "No distributors mapped to you", "Contact an Admin"),
    ("Review & Submit is disabled", "Already submitted this period", "Pick a different distributor/period, or ask an Admin"),
    ("Bad Stock card missing", "Zero YTD sell-through", "Expected — not an error"),
    ("Logged out unexpectedly", "You just changed your password", "Log back in with the new password"),
]
table_shape = slide.shapes.add_table(len(rows), 3, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.0))
table = table_shape.table
table.columns[0].width = Inches(3.6)
table.columns[1].width = Inches(3.8)
table.columns[2].width = Inches(4.7)
for r, row in enumerate(rows):
    for c, text in enumerate(row):
        cell = table.cell(r, c)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13 if r else 14)
            p.font.bold = (r == 0)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r == 0 else DARK_TEXT
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else (LIGHT if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF))
add_footer(slide, "Distributor Operational Assessment — User Training")

prs.save(str(ROOT / "User_Training.pptx"))
print("Saved User_Training.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
